"""
Generate Q3 HR Presentation PowerPoint from q3_hr_ppt.html content.
Dark-themed, professional slides matching the HTML design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
BG       = RGBColor(0x0B, 0x0D, 0x17)
BG2      = RGBColor(0x11, 0x13, 0x27)
CARD_BG  = RGBColor(0x16, 0x18, 0x2D)
PURPLE   = RGBColor(0x7C, 0x6A, 0xFF)
BLUE     = RGBColor(0x4D, 0xA8, 0xFF)
TEAL     = RGBColor(0x3E, 0xDD, 0xC6)
PINK     = RGBColor(0xFF, 0x6B, 0x9D)
GOLD     = RGBColor(0xFF, 0xB7, 0x4D)
GREEN    = RGBColor(0x66, 0xDE, 0x93)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT2    = RGBColor(0xB0, 0xB0, 0xC0)
TEXT3    = RGBColor(0x70, 0x70, 0x90)
RED_SOFT = RGBColor(0x40, 0x1A, 0x28)
GRN_SOFT = RGBColor(0x14, 0x3A, 0x30)
BANNER_BG = RGBColor(0x3A, 0x35, 0x7A)
BANNER_T  = RGBColor(0x24, 0x5A, 0x5A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # blank layout


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, w, h, text, font_size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_name='Calibri', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multiline(slide, left, top, w, h, lines, font_size=14, color=TEXT2,
                  spacing=Pt(6), font_name='Calibri', align=PP_ALIGN.LEFT, bold=False):
    """lines is a list of strings."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = align
        p.space_before = spacing
        p.space_after = Pt(2)
    return txBox


def add_card(slide, left, top, w, h, emoji, title, desc, border_color=None):
    add_rect(slide, left, top, w, h, CARD_BG, border_color or RGBColor(0x25, 0x28, 0x45), 0.04)
    add_text(slide, left + Inches(0.25), top + Inches(0.2), w - Inches(0.5), Inches(0.5),
             emoji, font_size=28, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.15), top + Inches(0.7), w - Inches(0.3), Inches(0.4),
             title, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(slide, left + Inches(0.15), top + Inches(1.1), w - Inches(0.3), Inches(0.8),
             desc, font_size=10, color=TEXT2, align=PP_ALIGN.CENTER)


def add_flow_step(slide, left, top, emoji, title, desc, show_arrow=True):
    """Single flow step box."""
    add_rect(slide, left, top, Inches(1.6), Inches(1.5), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.06)
    add_text(slide, left, top + Inches(0.1), Inches(1.6), Inches(0.45),
             emoji, font_size=24, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.05), top + Inches(0.55), Inches(1.5), Inches(0.35),
             title, font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.05), top + Inches(0.9), Inches(1.5), Inches(0.5),
             desc, font_size=8, color=TEXT3, align=PP_ALIGN.CENTER)
    if show_arrow:
        add_text(slide, left + Inches(1.6), top + Inches(0.4), Inches(0.4), Inches(0.4),
                 '→', font_size=18, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)


def add_banner_item(slide, left, top, w, num, label, bg_color):
    add_rect(slide, left, top, w, Inches(1.1), bg_color, radius=0.08)
    add_text(slide, left, top + Inches(0.1), w, Inches(0.5),
             num, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(slide, left, top + Inches(0.6), w, Inches(0.4),
             label, font_size=10, color=RGBColor(0xDD, 0xDD, 0xEE), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.4),
         'QUARTER 3 PERFORMANCE REVIEW', font_size=12, color=PURPLE, bold=True,
         align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.5),
         'Making Teams Work\nFaster & Smarter', font_size=44, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(s, Inches(2.5), Inches(3.4), Inches(8.3), Inches(0.8),
         'Delivered two automation projects that help development teams\nsave time, reduce repeated work, and solve problems faster.',
         font_size=16, color=TEXT2, align=PP_ALIGN.CENTER)
add_text(s, Inches(4), Inches(4.4), Inches(5.3), Inches(0.4),
         'Manprit Singh Panesar', font_size=14, color=TEXT3, align=PP_ALIGN.CENTER)

# Stat boxes
stat_data = [('2', 'Projects Delivered'), ('3+', 'Teams Benefited'), ('70%', 'Time Saved')]
stat_left = Inches(3.2)
for num, label in stat_data:
    add_rect(s, stat_left, Inches(5.2), Inches(2.1), Inches(1.4), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.06)
    add_text(s, stat_left, Inches(5.35), Inches(2.1), Inches(0.6),
             num, font_size=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name='Calibri')
    add_text(s, stat_left, Inches(5.95), Inches(2.1), Inches(0.4),
             label, font_size=10, color=TEXT3, align=PP_ALIGN.CENTER)
    stat_left += Inches(2.4)

# ═══════════════════════════════════════════
# SLIDE 2: WHAT I DELIVERED
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.5), Inches(4), Inches(0.3),
         '📋  QUARTER AT A GLANCE', font_size=11, color=PURPLE, bold=True)
add_text(s, Inches(0.6), Inches(1.0), Inches(10), Inches(0.6),
         'What I Delivered This Quarter', font_size=34, bold=True, color=WHITE, font_name='Calibri')
add_text(s, Inches(0.6), Inches(1.7), Inches(9), Inches(0.5),
         'Two projects focused on making teams more productive — by removing repetitive work and making information easier to access.',
         font_size=14, color=TEXT2)

# Project 1 card
add_rect(s, Inches(0.6), Inches(2.8), Inches(5.9), Inches(3.5), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.03)
add_text(s, Inches(1.0), Inches(3.0), Inches(5), Inches(0.5),
         '⚙️', font_size=36)
add_text(s, Inches(1.0), Inches(3.6), Inches(5.2), Inches(0.5),
         'Project 1: Standardized Build Process for React', font_size=16, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(4.2), Inches(5.2), Inches(1.5),
         'Created a one-click setup system so any React project team can start working with the same quality process — instead of each team creating their own from scratch.',
         font_size=13, color=TEXT2)

# Project 2 card
add_rect(s, Inches(6.8), Inches(2.8), Inches(5.9), Inches(3.5), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.03)
add_text(s, Inches(7.2), Inches(3.0), Inches(5), Inches(0.5),
         '🤖', font_size=36)
add_text(s, Inches(7.2), Inches(3.6), Inches(5.2), Inches(0.5),
         'Project 2: AI-Powered Problem Solver (Jenkins MCP)', font_size=16, bold=True, color=WHITE)
add_text(s, Inches(7.2), Inches(4.2), Inches(5.2), Inches(1.5),
         'Connected our build system (Jenkins) with AI so anyone in the team — developers, QA, or managers — can simply ask "What went wrong?" and get instant answers in plain language.',
         font_size=13, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 3: PROJECT 1 — BEFORE vs AFTER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)

# Header
add_rect(s, Inches(0.6), Inches(0.4), Inches(0.7), Inches(0.7), PURPLE, radius=0.1)
add_text(s, Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.65),
         '01', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(0.4), Inches(8), Inches(0.4),
         'React UnifiedCI — Standardized Build Process', font_size=22, bold=True, color=WHITE)
add_text(s, Inches(1.5), Inches(0.85), Inches(8), Inches(0.3),
         'Making sure all React project teams follow the same quality process', font_size=12, color=TEXT2)

add_text(s, Inches(0.6), Inches(1.6), Inches(10), Inches(0.5),
         'Every Team Was Doing It Differently', font_size=26, bold=True, color=WHITE)

# BEFORE box
add_rect(s, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.5), RED_SOFT, PINK, 0.03)
add_text(s, Inches(1.0), Inches(2.45), Inches(5), Inches(0.4),
         '❌  Before', font_size=18, bold=True, color=PINK)
before_items = [
    '📄  Every team created their own process from scratch',
    '🔧  Different approaches led to different quality levels',
    '🐛  A fix in one project didn\'t help other projects',
    '⚠️  Some teams skipped important quality checks',
    '⏱️  Setting up a new project took days of effort',
]
add_multiline(s, Inches(1.0), Inches(3.1), Inches(5.0), Inches(3.5),
              before_items, font_size=12, color=TEXT2, spacing=Pt(10))

# VS text
add_text(s, Inches(6.4), Inches(4.0), Inches(0.6), Inches(0.5),
         'VS', font_size=14, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)

# AFTER box
add_rect(s, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.5), GRN_SOFT, TEAL, 0.03)
add_text(s, Inches(7.3), Inches(2.45), Inches(5), Inches(0.4),
         '✅  After My Work', font_size=18, bold=True, color=TEAL)
after_items = [
    '📦  One ready-made system shared by all teams',
    '🎯  Every project follows the same quality standard',
    '🔄  Fix once — it helps every project automatically',
    '✅  Quality checks happen automatically, nothing skipped',
    '🚀  New project setup now takes minutes',
]
add_multiline(s, Inches(7.3), Inches(3.1), Inches(5.0), Inches(3.5),
              after_items, font_size=12, color=TEXT2, spacing=Pt(10))

# ═══════════════════════════════════════════
# SLIDE 4: PROJECT 1 — HOW IT WORKS + METRICS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '⚙️  REACT UNIFIEDCI', font_size=11, color=PURPLE, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'Simple, Automatic, No Extra Work', font_size=28, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.4), Inches(9), Inches(0.4),
         'Teams just connect their project — the system handles everything else automatically.',
         font_size=13, color=TEXT2)

# Flow steps
flow_data = [
    ('👨‍💻', 'Developer Saves Work', 'Code is submitted'),
    ('🔍', 'Quality Check', 'Checks for errors'),
    ('🧪', 'Build & Test', 'Runs all tests'),
    ('📊', 'Report Created', 'Pass/fail summary'),
    ('🔔', 'Team Notified', 'Email & Slack alerts'),
]
x_start = Inches(0.6)
for i, (emoji, title, desc) in enumerate(flow_data):
    add_flow_step(s, x_start + Inches(i * 2.2), Inches(2.2), emoji, title, desc, show_arrow=(i < 4))

# Banner metrics
metrics = [('90%', 'Faster Setup'), ('100%', 'Quality Automated'), ('1', 'System for All'), ('0', 'Manual Work')]
bw = Inches(2.8)
for i, (num, label) in enumerate(metrics):
    add_banner_item(s, Inches(0.6) + bw * i + Inches(i * 0.2), Inches(4.6), bw, num, label, BANNER_BG)

# Results cards
result_data = [
    ('📦', 'One Shared System', 'All React teams use the same process now'),
    ('✅', 'Auto Quality Checks', 'Errors caught early, nothing missed'),
    ('🧪', 'Automatic Testing', 'Tests run on their own with reports'),
    ('📧', 'Instant Notifications', 'Teams know the status right away'),
]
cw = Inches(2.85)
for i, (em, t, d) in enumerate(result_data):
    add_card(s, Inches(0.6) + i * (cw + Inches(0.2)), Inches(5.95), cw, Inches(1.3), em, t, d)

# ═══════════════════════════════════════════
# SLIDE 5: PROJECT 2 — PROBLEM
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)

add_rect(s, Inches(0.6), Inches(0.4), Inches(0.7), Inches(0.7), TEAL, radius=0.1)
add_text(s, Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.65),
         '02', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(0.4), Inches(8), Inches(0.4),
         'Jenkins MCP — AI-Powered Problem Solver', font_size=22, bold=True, color=WHITE)
add_text(s, Inches(1.5), Inches(0.85), Inches(8), Inches(0.3),
         'Making it easy for anyone to understand and fix issues', font_size=12, color=TEXT2)

add_text(s, Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
         'Finding & Fixing Problems Was Painful', font_size=26, bold=True, color=WHITE)

# Problem cards
problems = [
    ('⏰', 'Hours Wasted',
     'When something went wrong, someone had to manually go through thousands of lines of data to find what happened. This took hours every time.'),
    ('🧑‍💼', 'Only Experts Could Help',
     'Only a few experienced people could understand the error details. Everyone else had to wait — creating bottlenecks.'),
    ('🚧', 'Everyone Depended on One Team',
     'QA, managers, and other teams always had to ask the same team for updates — they couldn\'t check on their own.'),
]
for i, (em, t, d) in enumerate(problems):
    cx = Inches(0.6) + i * Inches(4.1)
    add_rect(s, cx, Inches(2.3), Inches(3.8), Inches(3.0), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.04)
    add_text(s, cx + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.5), em, font_size=32)
    add_text(s, cx + Inches(0.2), Inches(3.05), Inches(3.4), Inches(0.35),
             t, font_size=15, bold=True, color=WHITE)
    add_text(s, cx + Inches(0.2), Inches(3.5), Inches(3.4), Inches(1.5),
             d, font_size=11, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 6: PROJECT 2 — SOLUTION (BEFORE vs AFTER)
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '💡  JENKINS MCP — THE SOLUTION', font_size=11, color=TEAL, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'Just Ask — AI Does the Rest', font_size=28, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.4), Inches(9), Inches(0.4),
         'Now anyone can ask a question in plain English and get an instant answer. No expertise needed.',
         font_size=13, color=TEXT2)

# BEFORE
add_rect(s, Inches(0.6), Inches(2.1), Inches(5.8), Inches(3.5), RED_SOFT, PINK, 0.03)
add_text(s, Inches(1.0), Inches(2.25), Inches(5), Inches(0.4),
         '❌  The Old Way', font_size=17, bold=True, color=PINK)
old_items = [
    '🖥️  Log into the system',
    '🔎  Navigate to the right project',
    '📜  Open the error details',
    '👀  Read through everything manually',
    '🤔  Guess what went wrong & try again',
]
add_multiline(s, Inches(1.0), Inches(2.85), Inches(5.0), Inches(2.5),
              old_items, font_size=12, color=TEXT2, spacing=Pt(8))

add_text(s, Inches(6.4), Inches(3.5), Inches(0.6), Inches(0.5),
         'VS', font_size=14, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)

# AFTER
add_rect(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(3.5), GRN_SOFT, TEAL, 0.03)
add_text(s, Inches(7.3), Inches(2.25), Inches(5), Inches(0.4),
         '✅  The AI Way', font_size=17, bold=True, color=TEAL)
new_items = [
    '💬  Just ask: "What went wrong?"',
    '⚡  AI pulls up the relevant information instantly',
    '🧠  AI identifies the root cause',
    '🔧  AI suggests the exact solution',
    '✅  You approve — problem fixed',
]
add_multiline(s, Inches(7.3), Inches(2.85), Inches(5.0), Inches(2.5),
              new_items, font_size=12, color=TEXT2, spacing=Pt(8))

# Who benefits
add_text(s, Inches(0.6), Inches(5.9), Inches(10), Inches(0.4),
         'Helpful for Everyone — Not Just Experts', font_size=16, bold=True, color=WHITE)
benefit_data = [
    ('👨‍💻', 'Developers', '"Why did my project fail?"\n→ Instant answer + fix'),
    ('🧪', 'QA Team', '"Show all failures from today"\n→ Full summary in seconds'),
    ('👔', 'Managers', '"Is this ready for release?"\n→ AI checks & gives clear answer'),
]
for i, (em, t, d) in enumerate(benefit_data):
    cx = Inches(0.6) + i * Inches(4.1)
    add_card(s, cx, Inches(6.3), Inches(3.8), Inches(1.1), em, t, d)

# ═══════════════════════════════════════════
# SLIDE 7: PROJECT 2 — FLOW DIAGRAM
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '⚙️  JENKINS MCP', font_size=11, color=TEAL, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'How It Works — From Question to Solution', font_size=28, bold=True, color=WHITE)

# Flow steps
flow2 = [
    ('💬', 'Anyone Asks', 'Plain English question'),
    ('🧠', 'AI Understands', 'Figures out what you need'),
    ('📡', 'Gets the Data', 'Pulls relevant info'),
    ('🔍', 'Finds Root Cause', 'Identifies what went wrong'),
    ('✅', 'Gives Solution', 'Ready-to-apply fix'),
]
for i, (emoji, title, desc) in enumerate(flow2):
    add_flow_step(s, Inches(0.6) + Inches(i * 2.3), Inches(2.0), emoji, title, desc, show_arrow=(i < 4))

# Example scenarios
add_text(s, Inches(0.6), Inches(4.0), Inches(10), Inches(0.4),
         'Real-World Examples', font_size=18, bold=True, color=WHITE)

scenarios = [
    ('Developer asks:', '"My build failed for the React app — what happened?"',
     'AI: "The test in LoginComponent.test.js failed because the expected text \'Welcome\' was not found. Suggested fix: Update the component to render \'Welcome\' instead of \'Hello\'."'),
    ('QA asks:', '"How many builds failed today?"',
     'AI: "3 out of 12 builds failed today. Here\'s a summary of each failure with root causes and suggested fixes."'),
    ('Manager asks:', '"Is the latest release ready to deploy?"',
     'AI: "The latest build passed all quality checks and tests. It\'s ready for deployment. Last successful build: 2 hours ago."'),
]
for i, (role, q, a) in enumerate(scenarios):
    cy = Inches(4.5) + i * Inches(1.0)
    add_rect(s, Inches(0.6), cy, Inches(12.1), Inches(0.85), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.02)
    add_text(s, Inches(0.8), cy + Inches(0.05), Inches(1.5), Inches(0.3),
             role, font_size=10, bold=True, color=TEAL)
    add_text(s, Inches(0.8), cy + Inches(0.3), Inches(3.5), Inches(0.45),
             q, font_size=10, color=GOLD)
    add_text(s, Inches(4.5), cy + Inches(0.08), Inches(8), Inches(0.7),
             a, font_size=9, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 8: COMBINED IMPACT
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '📊  OVERALL IMPACT', font_size=11, color=GOLD, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'Quarter 3 Results', font_size=34, bold=True, color=WHITE)
add_text(s, Inches(0.6), Inches(1.5), Inches(9), Inches(0.4),
         'Both projects together are saving teams significant time, improving quality, and removing bottlenecks.',
         font_size=14, color=TEXT2)

# Banner
metrics2 = [('70%', 'Less Time\nFinding Problems'), ('90%', 'Faster New\nProject Setup'),
            ('100%', 'Consistent\nQuality'), ('All', 'Teams\nSelf-Sufficient')]
for i, (num, label) in enumerate(metrics2):
    add_banner_item(s, Inches(0.6) + i * Inches(3.15), Inches(2.3), Inches(2.9), num, label, BANNER_T)

# Impact cards
add_rect(s, Inches(0.6), Inches(3.9), Inches(5.9), Inches(1.6), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.03)
add_text(s, Inches(1.0), Inches(4.0), Inches(0.5), Inches(0.4), '🎯', font_size=24)
add_text(s, Inches(1.0), Inches(4.4), Inches(5.2), Inches(0.3),
         'Same High Standard Everywhere', font_size=14, bold=True, color=WHITE)
add_text(s, Inches(1.0), Inches(4.8), Inches(5.2), Inches(0.6),
         'Every React project now follows the same quality process. No more variations between teams — one improvement benefits everyone.',
         font_size=11, color=TEXT2)

add_rect(s, Inches(6.8), Inches(3.9), Inches(5.9), Inches(1.6), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.03)
add_text(s, Inches(7.2), Inches(4.0), Inches(0.5), Inches(0.4), '🙌', font_size=24)
add_text(s, Inches(7.2), Inches(4.4), Inches(5.2), Inches(0.3),
         'Teams Work Independently', font_size=14, bold=True, color=WHITE)
add_text(s, Inches(7.2), Inches(4.8), Inches(5.2), Inches(0.6),
         'With AI assistance, anyone can check project status and understand issues on their own — no waiting, no bottlenecks.',
         font_size=11, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 9: LEARNINGS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '📚  GROWTH', font_size=11, color=PURPLE, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'What I Learned This Quarter', font_size=34, bold=True, color=WHITE)

learnings = [
    ('🏗️', 'Building Scalable Solutions', 'Learned to design systems that\nserve multiple teams with one solution'),
    ('🤖', 'Working with AI', 'Gained hands-on experience connecting\nAI to existing company systems'),
    ('🧪', 'Quality & Automation', 'Deepened understanding of automated\nquality checks, testing, and reporting'),
    ('🎯', 'Problem-First Approach', 'Focused on understanding real team\npain points before building solutions'),
    ('🗣️', 'Cross-Team Communication', 'Improved ability to work with and\npresent to different stakeholders'),
    ('📐', 'End-to-End Ownership', 'Took ownership from idea to delivery —\nplanning, building, testing, presenting'),
]
for i, (em, t, d) in enumerate(learnings):
    row = i // 3
    col = i % 3
    cx = Inches(0.6) + col * Inches(4.1)
    cy = Inches(1.8) + row * Inches(2.5)
    add_card(s, cx, cy, Inches(3.8), Inches(2.1), em, t, d)

# ═══════════════════════════════════════════
# SLIDE 10: NEXT QUARTER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '🔮  LOOKING AHEAD', font_size=11, color=TEAL, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'Plans for Quarter 4', font_size=34, bold=True, color=WHITE)

q4_plans = [
    ('🌐', 'Expand to More Platforms', 'Extend the standardized process to support more project types beyond React.'),
    ('☁️', 'Scale for Larger Teams', 'Make the system handle higher workloads automatically during busy periods.'),
    ('🛡️', 'Add Security Checks', 'Include automatic safety and compliance checks in every project\'s process.'),
]
for i, (em, t, d) in enumerate(q4_plans):
    cx = Inches(0.6) + i * Inches(4.1)
    add_rect(s, cx, Inches(2.0), Inches(3.8), Inches(3.5), CARD_BG, RGBColor(0x25, 0x28, 0x45), 0.04)
    add_text(s, cx + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.5), em, font_size=40)
    add_text(s, cx + Inches(0.2), Inches(2.9), Inches(3.4), Inches(0.4),
             t, font_size=17, bold=True, color=WHITE)
    add_text(s, cx + Inches(0.2), Inches(3.45), Inches(3.4), Inches(1.5),
             d, font_size=13, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 11: THANK YOU
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
         'Thank You', font_size=60, bold=True, color=PURPLE, align=PP_ALIGN.CENTER, font_name='Calibri')
add_text(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
         'Happy to answer any questions', font_size=18, color=TEXT2, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.4),
         'Manprit Singh Panesar · Quarter 3 Review', font_size=14, color=TEXT3, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = r'c:\Users\dell\Downloads\React-UnifiedCI\Q3_Review_Manprit_Singh_Panesar.pptx'
prs.save(output_path)
print(f'\n✅ PowerPoint saved to: {output_path}')
print(f'   Total slides: {len(prs.slides)}')
