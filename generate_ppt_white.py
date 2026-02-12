"""
Generate Q3 HR Presentation PowerPoint — WHITE THEME
Based on q3_hr_presentation.html (8-slide compact version).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── WHITE THEME Colors ──
BG       = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG  = RGBColor(0xF5, 0xF5, 0xF9)
CARD_BD  = RGBColor(0xE0, 0xE0, 0xE8)
PURPLE   = RGBColor(0x5B, 0x4C, 0xDB)
BLUE     = RGBColor(0x2E, 0x86, 0xDE)
TEAL     = RGBColor(0x0F, 0xA3, 0x8E)
PINK     = RGBColor(0xE0, 0x40, 0x6E)
GOLD     = RGBColor(0xD4, 0x8A, 0x20)
GREEN    = RGBColor(0x28, 0xA7, 0x5B)
BLACK    = RGBColor(0x1A, 0x1A, 0x2E)
TEXT1    = RGBColor(0x22, 0x22, 0x33)
TEXT2    = RGBColor(0x55, 0x55, 0x70)
TEXT3    = RGBColor(0x88, 0x88, 0x9A)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED_SOFT = RGBColor(0xFD, 0xED, 0xF0)
GRN_SOFT = RGBColor(0xE8, 0xF8, 0xF0)
BANNER_P = RGBColor(0x5B, 0x4C, 0xDB)  # Purple banner
BANNER_T = RGBColor(0x0F, 0xA3, 0x8E)  # Teal banner

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_text(slide, left, top, w, h, text, font_size=18, color=TEXT1, bold=False,
             align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return txBox


def add_multiline(slide, left, top, w, h, lines, font_size=14, color=TEXT2,
                  spacing=Pt(6), font_name='Calibri', align=PP_ALIGN.LEFT, bold=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = align
        p.space_before = spacing
        p.space_after = Pt(2)
    return txBox


def add_card(slide, left, top, w, h, emoji, title, desc, bg=None, bd=None):
    add_rect(slide, left, top, w, h, bg or CARD_BG, bd or CARD_BD, 0.04)
    add_text(slide, left + Inches(0.2), top + Inches(0.15), w - Inches(0.4), Inches(0.45),
             emoji, font_size=26, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.12), top + Inches(0.6), w - Inches(0.24), Inches(0.4),
             title, font_size=12, bold=True, color=TEXT1, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.12), top + Inches(0.95), w - Inches(0.24), Inches(0.8),
             desc, font_size=10, color=TEXT2, align=PP_ALIGN.CENTER)


def add_flow_step(slide, left, top, emoji, title, desc, show_arrow=True):
    add_rect(slide, left, top, Inches(1.6), Inches(1.5), CARD_BG, CARD_BD, 0.06)
    add_text(slide, left, top + Inches(0.1), Inches(1.6), Inches(0.4),
             emoji, font_size=22, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.05), top + Inches(0.5), Inches(1.5), Inches(0.35),
             title, font_size=10, bold=True, color=TEXT1, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.05), top + Inches(0.85), Inches(1.5), Inches(0.5),
             desc, font_size=8, color=TEXT3, align=PP_ALIGN.CENTER)
    if show_arrow:
        add_text(slide, left + Inches(1.6), top + Inches(0.35), Inches(0.45), Inches(0.4),
                 '→', font_size=18, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)


def add_banner_item(slide, left, top, w, num, label, bg_color):
    add_rect(slide, left, top, w, Inches(1.1), bg_color, radius=0.08)
    add_text(slide, left, top + Inches(0.1), w, Inches(0.5),
             num, font_size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.6), w, Inches(0.4),
             label, font_size=10, color=RGBColor(0xE8, 0xE8, 0xF0), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.5), Inches(0.8), Inches(12.3), Inches(0.4),
         'QUARTER 3 PERFORMANCE REVIEW', font_size=12, color=PURPLE, bold=True,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(1.6), Inches(11.3), Inches(1.5),
         'Making Teams Work\nFaster & Smarter', font_size=44, bold=True, color=TEXT1,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(2.5), Inches(3.4), Inches(8.3), Inches(0.8),
         'Delivered two automation projects that help development teams\nsave time, reduce repeated work, and solve problems faster.',
         font_size=16, color=TEXT2, align=PP_ALIGN.CENTER)
add_text(s, Inches(4), Inches(4.4), Inches(5.3), Inches(0.4),
         'Manprit Singh Panesar', font_size=14, color=TEXT3, align=PP_ALIGN.CENTER)

# Stat boxes
stat_data = [('2', 'Projects Delivered'), ('3+', 'Teams Benefited'), ('70%', 'Time Saved')]
stat_left = Inches(3.2)
for num, label in stat_data:
    add_rect(s, stat_left, Inches(5.2), Inches(2.1), Inches(1.4), CARD_BG, CARD_BD, 0.06)
    add_text(s, stat_left, Inches(5.35), Inches(2.1), Inches(0.6),
             num, font_size=36, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    add_text(s, stat_left, Inches(5.95), Inches(2.1), Inches(0.4),
             label, font_size=10, color=TEXT3, align=PP_ALIGN.CENTER)
    stat_left += Inches(2.4)

# ═══════════════════════════════════════════
# SLIDE 2: WHAT I DELIVERED
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.5), Inches(4), Inches(0.3),
         '📋  QUARTER AT A GLANCE', font_size=11, color=PURPLE, bold=True)
add_text(s, Inches(0.6), Inches(1.0), Inches(10), Inches(0.6),
         'What I Delivered This Quarter', font_size=34, bold=True, color=TEXT1)
add_text(s, Inches(0.6), Inches(1.7), Inches(9), Inches(0.5),
         'Focused on two projects — one to standardize how teams set up their projects, and one to let AI help everyone solve problems instantly.',
         font_size=14, color=TEXT2)

# Project 1 card
add_rect(s, Inches(0.6), Inches(2.8), Inches(5.9), Inches(3.5), CARD_BG, CARD_BD, 0.03)
add_text(s, Inches(1.0), Inches(3.0), Inches(5), Inches(0.5), '⚙️', font_size=36)
add_text(s, Inches(1.0), Inches(3.6), Inches(5.2), Inches(0.5),
         'React UnifiedCI — Standardized Project Setup', font_size=16, bold=True, color=TEXT1)
add_text(s, Inches(1.0), Inches(4.2), Inches(5.2), Inches(1.5),
         'Built a shared, ready-to-use system so every React web application team follows the same quality process — instead of each team creating their own setup from scratch every time.',
         font_size=13, color=TEXT2)

# Project 2 card
add_rect(s, Inches(6.8), Inches(2.8), Inches(5.9), Inches(3.5), CARD_BG, CARD_BD, 0.03)
add_text(s, Inches(7.2), Inches(3.0), Inches(5), Inches(0.5), '🤖', font_size=36)
add_text(s, Inches(7.2), Inches(3.6), Inches(5.2), Inches(0.5),
         'Jenkins MCP — AI-Powered Problem Solver', font_size=16, bold=True, color=TEXT1)
add_text(s, Inches(7.2), Inches(4.2), Inches(5.2), Inches(1.5),
         'Connected our build system (Jenkins) with AI, so anyone can simply ask "What went wrong?" in plain English and get an instant answer — no expertise needed.',
         font_size=13, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 3: REACT UNIFIEDCI — BEFORE vs AFTER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)

add_rect(s, Inches(0.6), Inches(0.4), Inches(0.7), Inches(0.7), PURPLE, radius=0.1)
add_text(s, Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.65),
         '01', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(0.4), Inches(8), Inches(0.4),
         'React UnifiedCI', font_size=22, bold=True, color=TEXT1)
add_text(s, Inches(1.5), Inches(0.85), Inches(8), Inches(0.3),
         'One shared system for all React web application projects', font_size=12, color=TEXT2)

# BEFORE box
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(4.2), RED_SOFT, PINK, 0.03)
add_text(s, Inches(1.0), Inches(1.85), Inches(5), Inches(0.4),
         '❌  Before', font_size=17, bold=True, color=PINK)
before_items = [
    '📄  Each team wrote their own setup from scratch',
    '🔧  Different quality levels across teams',
    '⚠️  Some teams skipped important checks',
    '⏱️  New project setup took days',
]
add_multiline(s, Inches(1.0), Inches(2.5), Inches(5.0), Inches(3.0),
              before_items, font_size=12, color=TEXT2, spacing=Pt(10))

add_text(s, Inches(6.4), Inches(3.5), Inches(0.6), Inches(0.5),
         'VS', font_size=14, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)

# AFTER box
add_rect(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.2), GRN_SOFT, TEAL, 0.03)
add_text(s, Inches(7.3), Inches(1.85), Inches(5), Inches(0.4),
         '✅  After', font_size=17, bold=True, color=TEAL)
after_items = [
    '📦  One shared system used by all teams',
    '🎯  Same quality standard everywhere',
    '✅  All checks run automatically, nothing skipped',
    '🚀  New project setup in minutes',
]
add_multiline(s, Inches(7.3), Inches(2.5), Inches(5.0), Inches(3.0),
              after_items, font_size=12, color=TEXT2, spacing=Pt(10))

# Banner
metrics = [('90%', 'Faster Setup'), ('100%', 'Quality Automated'), ('1', 'System for All'), ('0', 'Manual Work')]
bw = Inches(2.8)
for i, (num, label) in enumerate(metrics):
    add_banner_item(s, Inches(0.6) + i * (bw + Inches(0.2)), Inches(6.2), bw, num, label, BANNER_P)

# ═══════════════════════════════════════════
# SLIDE 4: REACT UNIFIEDCI — HOW IT WORKS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '⚙️  REACT UNIFIEDCI', font_size=11, color=PURPLE, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'How It Works — Fully Automatic', font_size=28, bold=True, color=TEXT1)
add_text(s, Inches(0.6), Inches(1.4), Inches(9), Inches(0.4),
         'Teams just connect their project — the system does the quality checking, testing, and reporting on its own.',
         font_size=13, color=TEXT2)

# Flow
flow_data = [
    ('👨‍💻', 'Developer Saves Work', 'Code is submitted'),
    ('🔍', 'Quality Check', 'Checks for errors automatically'),
    ('🧪', 'Build & Test', 'Builds project & runs tests'),
    ('📊', 'Report Created', 'Clear pass/fail report'),
    ('🔔', 'Team Notified', 'Email & Slack alerts'),
]
for i, (emoji, title, desc) in enumerate(flow_data):
    add_flow_step(s, Inches(0.6) + Inches(i * 2.2), Inches(2.3), emoji, title, desc, i < 4)

# Result cards
results = [
    ('📦', 'One Shared System', 'All React teams use the same process now'),
    ('✅', 'Auto Quality Checks', 'Errors caught early, nothing missed'),
    ('🧪', 'Automatic Testing', 'Tests run on their own with clear reports'),
    ('📧', 'Instant Notifications', 'Teams know the status right away'),
]
cw = Inches(2.85)
for i, (em, t, d) in enumerate(results):
    add_card(s, Inches(0.6) + i * (cw + Inches(0.2)), Inches(4.5), cw, Inches(1.6), em, t, d)

# ═══════════════════════════════════════════
# SLIDE 5: JENKINS MCP — BEFORE vs AFTER
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)

add_rect(s, Inches(0.6), Inches(0.4), Inches(0.7), Inches(0.7), TEAL, radius=0.1)
add_text(s, Inches(0.6), Inches(0.45), Inches(0.7), Inches(0.65),
         '02', font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(0.4), Inches(8), Inches(0.4),
         'Jenkins MCP — AI-Powered Problem Solver', font_size=22, bold=True, color=TEXT1)
add_text(s, Inches(1.5), Inches(0.85), Inches(8), Inches(0.3),
         'Anyone can ask questions in plain English and get instant answers', font_size=12, color=TEXT2)

# BEFORE
add_rect(s, Inches(0.6), Inches(1.7), Inches(5.8), Inches(3.5), RED_SOFT, PINK, 0.03)
add_text(s, Inches(1.0), Inches(1.85), Inches(5), Inches(0.4),
         '❌  The Old Way', font_size=17, bold=True, color=PINK)
old_items = [
    '🖥️  Log into the system manually',
    '📜  Scroll through thousands of lines to find the error',
    '🤔  Only experts could understand the issues',
    '⏰  Took hours to find and fix a single problem',
]
add_multiline(s, Inches(1.0), Inches(2.5), Inches(5.0), Inches(2.5),
              old_items, font_size=12, color=TEXT2, spacing=Pt(10))

add_text(s, Inches(6.4), Inches(3.1), Inches(0.6), Inches(0.5),
         'VS', font_size=14, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)

# AFTER
add_rect(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(3.5), GRN_SOFT, TEAL, 0.03)
add_text(s, Inches(7.3), Inches(1.85), Inches(5), Inches(0.4),
         '✅  The AI Way', font_size=17, bold=True, color=TEAL)
new_items = [
    '💬  Just ask: "What went wrong?"',
    '⚡  AI finds the issue instantly',
    '🧠  Anyone can use it — no expertise needed',
    '✅  Get the fix suggestion in seconds',
]
add_multiline(s, Inches(7.3), Inches(2.5), Inches(5.0), Inches(2.5),
              new_items, font_size=12, color=TEXT2, spacing=Pt(10))

# Who benefits
benefit_data = [
    ('👨‍💻', 'Developers', '"Why did my project fail?"\n→ Instant answer + fix'),
    ('🧪', 'QA Team', '"Show all failures from today"\n→ Full summary in seconds'),
    ('👔', 'Managers', '"Is this ready for release?"\n→ AI checks & gives clear answer'),
]
for i, (em, t, d) in enumerate(benefit_data):
    add_card(s, Inches(0.6) + i * Inches(4.1), Inches(5.6), Inches(3.8), Inches(1.6), em, t, d)

# ═══════════════════════════════════════════
# SLIDE 6: JENKINS MCP — HOW IT WORKS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '⚙️  JENKINS MCP', font_size=11, color=TEAL, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'How It Works — Ask and Get Answers', font_size=28, bold=True, color=TEXT1)

# Flow
flow2 = [
    ('💬', 'Anyone Asks', 'Plain English question'),
    ('🧠', 'AI Understands', 'Figures out what you need'),
    ('📡', 'Fetches Info', 'Gets the right data'),
    ('🔍', 'Finds the Answer', 'What went wrong & why'),
    ('✅', 'Gives Solution', 'Ready-to-apply fix'),
]
for i, (emoji, title, desc) in enumerate(flow2):
    add_flow_step(s, Inches(0.6) + Inches(i * 2.3), Inches(2.0), emoji, title, desc, i < 4)

# Example scenarios
add_text(s, Inches(0.6), Inches(4.0), Inches(10), Inches(0.4),
         'Real-World Examples', font_size=18, bold=True, color=TEXT1)

scenarios = [
    ('Developer asks:', '"My build failed — what happened?"',
     'AI: "The test failed because the expected text was not found. Suggested fix: Update the component text."'),
    ('QA asks:', '"How many builds failed today?"',
     'AI: "3 out of 12 builds failed today. Here\'s a summary of each failure with root causes and fixes."'),
    ('Manager asks:', '"Is the latest release ready to deploy?"',
     'AI: "The latest build passed all checks and tests. It\'s ready for deployment."'),
]
for i, (role, q, a) in enumerate(scenarios):
    cy = Inches(4.6) + i * Inches(0.95)
    add_rect(s, Inches(0.6), cy, Inches(12.1), Inches(0.8), CARD_BG, CARD_BD, 0.02)
    add_text(s, Inches(0.8), cy + Inches(0.05), Inches(1.5), Inches(0.3),
             role, font_size=10, bold=True, color=TEAL)
    add_text(s, Inches(0.8), cy + Inches(0.3), Inches(3.2), Inches(0.4),
             q, font_size=10, color=GOLD)
    add_text(s, Inches(4.2), cy + Inches(0.08), Inches(8.2), Inches(0.6),
             a, font_size=9, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 7: IMPACT + LEARNINGS
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(0.6), Inches(0.4), Inches(4), Inches(0.3),
         '📊  RESULTS & GROWTH', font_size=11, color=GOLD, bold=True)
add_text(s, Inches(0.6), Inches(0.85), Inches(10), Inches(0.5),
         'Quarter 3 Impact & Learnings', font_size=34, bold=True, color=TEXT1)

# Banner
metrics2 = [('70%', 'Less Time\nSolving Problems'), ('90%', 'Faster\nProject Setup'),
            ('100%', 'Consistent\nQuality'), ('All', 'Teams\nSelf-Sufficient')]
for i, (num, label) in enumerate(metrics2):
    add_banner_item(s, Inches(0.6) + i * Inches(3.15), Inches(1.7), Inches(2.9), num, label, BANNER_T)

# Learnings
add_text(s, Inches(0.6), Inches(3.2), Inches(10), Inches(0.4),
         'Key Learnings', font_size=16, bold=True, color=TEXT1)

learnings = [
    ('🏗️', 'Building Scalable Solutions', 'Designing one system that\nserves multiple teams'),
    ('🤖', 'Working with AI', 'Connecting AI to existing\ncompany systems'),
    ('🎯', 'End-to-End Ownership', 'From idea to delivery —\nplanning, building, presenting'),
]
for i, (em, t, d) in enumerate(learnings):
    add_card(s, Inches(0.6) + i * Inches(4.1), Inches(3.7), Inches(3.8), Inches(1.8), em, t, d)

# Impact cards
add_rect(s, Inches(0.6), Inches(5.8), Inches(5.9), Inches(1.4), CARD_BG, CARD_BD, 0.03)
add_text(s, Inches(1.0), Inches(5.9), Inches(0.5), Inches(0.4), '🎯', font_size=22)
add_text(s, Inches(1.0), Inches(6.25), Inches(5.2), Inches(0.25),
         'Same High Standard Everywhere', font_size=13, bold=True, color=TEXT1)
add_text(s, Inches(1.0), Inches(6.55), Inches(5.2), Inches(0.5),
         'Every project follows the same quality process — one improvement benefits everyone.',
         font_size=10, color=TEXT2)

add_rect(s, Inches(6.8), Inches(5.8), Inches(5.9), Inches(1.4), CARD_BG, CARD_BD, 0.03)
add_text(s, Inches(7.2), Inches(5.9), Inches(0.5), Inches(0.4), '🙌', font_size=22)
add_text(s, Inches(7.2), Inches(6.25), Inches(5.2), Inches(0.25),
         'Teams Work Independently', font_size=13, bold=True, color=TEXT1)
add_text(s, Inches(7.2), Inches(6.55), Inches(5.2), Inches(0.5),
         'Anyone can check project status and understand issues on their own — no bottlenecks.',
         font_size=10, color=TEXT2)

# ═══════════════════════════════════════════
# SLIDE 8: THANK YOU
# ═══════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG)
add_text(s, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
         'Thank You', font_size=60, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.8), Inches(11.3), Inches(0.5),
         'Happy to answer any questions', font_size=18, color=TEXT2, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.4),
         'Manprit Singh Panesar · Quarter 3 Review', font_size=14, color=TEXT3, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = r'c:\Users\dell\Downloads\React-UnifiedCI\Q3_Review_White_Theme.pptx'
prs.save(output_path)
print(f'\n✅ PowerPoint saved to: {output_path}')
print(f'   Total slides: {len(prs.slides)}')
